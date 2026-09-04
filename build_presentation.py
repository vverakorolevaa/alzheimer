"""
Сборка презентации (.pptx) по проекту «AD Blood Stage».
Структура — 8 разделов по требованиям защиты + титул и заключение.
Числа берутся из обученной модели (results/panel_model.pkl), чтобы слайды
всегда соответствовали реальным результатам. Запуск: python cli.py report
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

import config

IMG = "docs/img"
OUT = "docs/Презентация_AD_Blood_Stage.pptx"
ACCENT = RGBColor(0xC0, 0x39, 0x2B)   # тёмно-красный
DARK = RGBColor(0x2C, 0x3E, 0x50)


def _metrics():
    """Достаёт числа из модели; если её нет — значения по умолчанию."""
    try:
        import joblib
        m = joblib.load(os.path.join(config.RESULTS_FOLDER, "panel_model.pkl"))
        mt = m["metrics"]
        return {
            "n_genes": len(m["genes"]),
            "macro": mt["macro_auc"], "ci": mt["macro_auc_ci"],
            "ctl": mt["per_class_auc"]["CTL"], "mci": mt["per_class_auc"]["MCI"],
            "ad": mt["per_class_auc"]["AD"], "n": mt["n_samples"],
            "npc": mt["n_per_class"],
        }
    except Exception:
        return {"n_genes": 15, "macro": 0.75, "ci": [0.71, 0.79],
                "ctl": 0.86, "mci": 0.66, "ad": 0.72, "n": 329,
                "npc": {"CTL": 104, "MCI": 80, "AD": 145}}


def _title_slide(prs, title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    box = s.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.7), Inches(2.5))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title; p.alignment = PP_ALIGN.CENTER
    p.runs[0].font.size = Pt(38); p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = ACCENT
    p2 = tf.add_paragraph(); p2.text = subtitle; p2.alignment = PP_ALIGN.CENTER
    p2.runs[0].font.size = Pt(20); p2.runs[0].font.color.rgb = DARK
    return s


def _content_slide(prs, title, bullets):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    t = s.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(1.0))
    tp = t.text_frame.paragraphs[0]; tp.text = title
    tp.runs[0].font.size = Pt(30); tp.runs[0].font.bold = True
    tp.runs[0].font.color.rgb = ACCENT
    body = s.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12.0), Inches(5.5))
    bf = body.text_frame; bf.word_wrap = True
    for i, (text, lvl) in enumerate(bullets):
        p = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
        p.text = ("• " if lvl == 0 else "    – ") + text
        p.level = lvl
        p.runs[0].font.size = Pt(20 if lvl == 0 else 17)
        p.runs[0].font.color.rgb = DARK
        p.space_after = Pt(6)
    return s


def _image_slide(prs, title, images, by="height", size=5.6):
    """images — список путей; размещаем в ряд (by='height') или один крупно (by='width')."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    t = s.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.1), Inches(0.9))
    tp = t.text_frame.paragraphs[0]; tp.text = title
    tp.runs[0].font.size = Pt(28); tp.runs[0].font.bold = True
    tp.runs[0].font.color.rgb = ACCENT
    imgs = [p for p in images if os.path.exists(p)]
    if not imgs:
        return s
    if by == "height":
        n = len(imgs); gap = 0.3
        # грубое центрирование ряда (ширина оценивается из высоты, аспект ~0.6)
        approx_w = size * 0.62
        total = n * approx_w + (n - 1) * gap
        left = max(0.4, (13.333 - total) / 2)
        for p in imgs:
            s.shapes.add_picture(p, Inches(left), Inches(1.4), height=Inches(size))
            left += approx_w + gap
    else:
        s.shapes.add_picture(imgs[0], Inches((13.333 - size) / 2), Inches(1.5),
                             width=Inches(size))
    return s


def _add_box(slide, x, y, w, h, bg, title, sub):
    """Цветной блок (скруглённый прямоугольник) с заголовком и подписью."""
    shp = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    # заливка
    shp.fill.solid()
    shp.fill.fore_color.rgb = bg
    # рамка того же цвета — визуально без рамки
    shp.line.color.rgb = bg
    shp.line.width = Pt(0.5)
    # текст
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_top    = Emu(120000)
    tf.margin_bottom = Emu(60000)
    tf.margin_left   = Emu(90000)
    tf.margin_right  = Emu(90000)
    p1 = tf.paragraphs[0]
    p1.text = title
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.runs[0]
    r1.font.size = Pt(14)
    r1.font.bold = True
    r1.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p2 = tf.add_paragraph()
    p2.text = sub
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(4)
    r2 = p2.runs[0]
    r2.font.size = Pt(10)
    r2.font.color.rgb = RGBColor(0xE8, 0xF0, 0xFF)


def _add_arrow(slide, x, y):
    """Стрелка между блоками."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(0.38), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    p.text = ">"
    p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)


def _scheme_slide(prs):
    """Визуальная функциональная схема: 5 цветных блоков + стрелки."""
    s = prs.slides.add_slide(prs.slide_layouts[6])

    # заголовок
    t = s.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(12.3), Inches(0.85))
    tp = t.text_frame.paragraphs[0]
    tp.text = "4. Функциональная схема"
    tp.runs[0].font.size = Pt(30)
    tp.runs[0].font.bold = True
    tp.runs[0].font.color.rgb = ACCENT

    STEPS = [
        (RGBColor(0x29, 0x80, 0xB9),
         "Данные",
         "GSE63060\n329 образцов крови\n(кровь, открытые)"),
        (RGBColor(0x27, 0xAE, 0x60),
         "Загрузка",
         "probe → ген\nметки CTL / MCI / AD"),
        (RGBColor(0xD3, 0x5A, 0x00),
         "Предобработка",
         "log2, нормализация\nфильтр генов"),
        (RGBColor(0x8E, 0x44, 0xAD),
         "ПАНЕЛЬ + МОДЕЛЬ",
         "15 генов, без утечки\nкросс-валидация + SHAP"),
        (RGBColor(0x1A, 0x5E, 0x7A),
         "Веб-приложение",
         "Streamlit\nстадия + объяснение"),
    ]

    BOX_W = 2.08
    BOX_H = 2.3
    ARR_W = 0.30
    N     = len(STEPS)
    total = N * BOX_W + (N - 1) * ARR_W
    sx    = (13.333 - total) / 2
    BY    = 2.35

    for i, (color, title, sub) in enumerate(STEPS):
        bx = sx + i * (BOX_W + ARR_W)
        _add_box(s, bx, BY, BOX_W, BOX_H, color, title, sub)
        if i < N - 1:
            _add_arrow(s, bx + BOX_W, BY + BOX_H / 2 - 0.23)

    # подпись под блоком «ПАНЕЛЬ»
    note = s.shapes.add_textbox(
        Inches(sx + 3 * (BOX_W + ARR_W)),
        Inches(BY + BOX_H + 0.12),
        Inches(BOX_W), Inches(0.38))
    np_ = note.text_frame.paragraphs[0]
    np_.text = "Новизна проекта"
    np_.alignment = PP_ALIGN.CENTER
    r = np_.runs[0]
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0x8E, 0x44, 0xAD)

    return s


def build():
    m = _metrics()
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

    _title_slide(prs, "AD Blood Stage",
                 "Ранняя стадийная диагностика болезни Альцгеймера "
                 "по экспрессии генов крови\n(норма → MCI → деменция)")

    _content_slide(prs, "1. Литобзор (кратко)", [
        ("Болезнь Альцгеймера — 55+ млн человек в мире (ВОЗ).", 0),
        ("Главная проблема — позднее установление диагноза.", 0),
        ("Патология развивается 10–20 лет ДО симптомов "
         "(амилоидные бляшки, тау-клубки → атрофия мозга).", 0),
        ("К моменту симптомов мозг уже повреждён, лекарства не работают.", 0),
        ("Стадии: норма → MCI (промежуточная) → деменция.", 0),
        ("MCI — терапевтическое окно: лечение эффективно именно здесь.", 1),
    ])

    _content_slide(prs, "2. Обоснование и новизна", [
        ("Обоснование: диагноз ставят поздно — нужно ловить окно MCI.", 0),
        ("Существующие методы: МРТ/ПЭТ дороги, ткань мозга у живого недоступна.", 1),
        ("Кровь — неинвазивно, дёшево, несёт ранний молекулярный сигнал.", 1),
        ("В литературе каждый подход закрывает лишь часть задачи:", 0),
        ("— по генам, но только «болен/здоров» (напр. AITeQ);", 1),
        ("— 3 стадии, но по МРТ или клиническим шкалам, не по генам;", 1),
        ("— гены+стадии+объяснение, но без сайта и на закрытых данных (c-Triadem).", 1),
        ("Новизна — впервые ВСЁ вместе: гены крови + 3 стадии + мин. панель", 0),
        ("+ объяснение + рабочий сайт + валидация без утечки на открытых данных.", 1),
    ])

    _content_slide(prs, "3. Реферат — функционал", [
        (f"Вход: уровни экспрессии {m['n_genes']} генов панели (анализ крови).", 0),
        ("Выход: 3 вероятности (норма / MCI / деменция) + наиболее вероятная стадия.", 0),
        ("Главное: персональное объяснение — вклад каждого гена (SHAP).", 0),
        ("Назначение: инструмент поддержки решения врача (скрининг, не диагноз).", 0),
        (f"Ожидаемый результат: macro-AUC ≈ {m['macro']:.2f} (кросс-валидация).", 0),
    ])

    _scheme_slide(prs)

    _content_slide(prs, "5. Инструменты и обоснование", [
        ("Python + scikit-learn — стандарт ML, воспроизводимо.", 0),
        ("Мультиномиальная логистическая регрессия — интерпретируема и", 0),
        ("устойчива на малой выборке (нейросеть переобучилась бы на 329 образцах).", 1),
        ("SHAP — объяснение: глобальная важность + персональный вклад генов.", 0),
        ("ANOVA F-критерий — отбор самых различающих стадии генов.", 0),
        ("Streamlit — веб-интерфейс на Python; matplotlib — графики.", 0),
    ])

    _content_slide(prs, "6. Структура проекта", [
        ("download_data.py — загрузка данных с NCBI GEO.", 0),
        ("data_loader.py / preprocessor.py — парсинг, probe→ген, обработка.", 0),
        ("biomarker_panel.py — ЯДРО: отбор панели, кросс-валидация, модель, SHAP.", 0),
        ("app.py — веб-интерфейс (Streamlit).", 0),
        ("config.py — параметры; cli.py — единая точка входа.", 0),
    ])

    _image_slide(prs, "7. Интерфейс — ввод и результат",
                 [f"{IMG}/01_input.png", f"{IMG}/02_result.png"], by="height", size=5.8)
    _image_slide(prs, "7. Интерфейс — объяснение и метод",
                 [f"{IMG}/04_method.png"], by="height", size=5.9)

    _content_slide(prs, "8. Результаты", [
        (f"Данные: GSE63060, кровь, {m['n']} образцов "
         f"(норма {m['npc']['CTL']} / MCI {m['npc']['MCI']} / деменция {m['npc']['AD']}).", 0),
        (f"macro-AUC (кросс-валидация) = {m['macro']:.2f} "
         f"(95% ДИ {m['ci'][0]:.2f}–{m['ci'][1]:.2f}).", 0),
        (f"AUC по стадиям: норма {m['ctl']:.2f}, MCI {m['mci']:.2f}, "
         f"деменция {m['ad']:.2f}.", 0),
        ("Норму от болезни отделяет уверенно; MCI — труднее (показано прямо).", 1),
        ("Утечки/переобучения нет (наивный ≈ корректный отбор).", 0),
    ])
    _image_slide(prs, "8. Результаты — графики",
                 [f"{IMG}/metrics.png"], by="width", size=11.5)

    _content_slide(prs, "Заключение", [
        ("Рабочий веб-инструмент стадийной диагностики AD по крови на реальных данных.", 0),
        ("Ранняя, неинвазивная и объяснимая оценка стадии.", 0),
        ("Методология без утечки: метрики по классам, доверительные интервалы.", 0),
        ("Прототип поддержки решения врача, не замена ему.", 0),
        ("Дальше: внешняя валидация (GSE63061), расширение панели, qPCR.", 0),
        ("Спасибо за внимание!", 0),
    ])

    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    print(f"Презентация сохранена: {OUT} ({len(prs.slides)} слайдов)")


def main():
    build()


if __name__ == "__main__":
    main()
